import os
import shutil
import subprocess
import tempfile
import warnings

from backend.paths import long_path

warnings.simplefilter(action="ignore", category=UserWarning)

SUPPORTED_EXTRACT = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".txt", ".md", ".py", ".java", ".html", ".csv"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg"}
VIDEO_EXTS = {".mp4", ".mov", ".wmv", ".m4v"}
TRANSCRIBE_EXTS = AUDIO_EXTS | VIDEO_EXTS


def _libreoffice_convert(src_path: str, target_ext: str) -> str | None:
    """Convert via LibreOffice headless if available."""
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo:
        return None
    tmp_dir = tempfile.mkdtemp()
    try:
        result = subprocess.run(
            [lo, "--headless", "--convert-to", target_ext, "--outdir", tmp_dir, src_path],
            capture_output=True, timeout=120,
        )
        if result.returncode == 0:
            base = os.path.splitext(os.path.basename(src_path))[0]
            out = os.path.join(tmp_dir, f"{base}.{target_ext}")
            if os.path.exists(out):
                return out
    except Exception:
        pass
    return None


def _clean_ole_text(text: str) -> str:
    import re
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text).replace("\r", "\n")
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _extract_doc_ole(path: str) -> tuple[str, bool]:
    """Pure-Python binary .doc extraction via the MS-DOC piece table (CLX/PlcPcd).

    Reconstructs the document text by walking the piece table in the table stream,
    decoding each piece as UTF-16LE or 8-bit (cp1255 — Hebrew, ASCII-compatible).
    Handles both contiguous and fragmented (fComplex) documents.
    """
    import struct
    try:
        import olefile
    except ImportError:
        return "[.doc: olefile not installed]", False
    try:
        ole = olefile.OleFileIO(path)
    except Exception as e:
        return f"[.doc: not a valid OLE file: {e}]", False

    try:
        if not ole.exists("WordDocument"):
            ole.close()
            return "[.doc: no WordDocument stream]", False

        wd = ole.openstream("WordDocument").read()
        flags = struct.unpack_from("<H", wd, 0x0A)[0]
        use_1table = bool(flags & 0x0200)        # fWhichTblStm
        fc_clx = struct.unpack_from("<I", wd, 0x01A2)[0]
        lcb_clx = struct.unpack_from("<I", wd, 0x01A6)[0]

        tbl_name = "1Table" if use_1table else "0Table"
        if not ole.exists(tbl_name):
            tbl_name = "0Table" if use_1table else "1Table"
        table = ole.openstream(tbl_name).read() if ole.exists(tbl_name) else b""
        ole.close()

        text = ""
        if table and lcb_clx > 0 and fc_clx + lcb_clx <= len(table):
            clx = table[fc_clx: fc_clx + lcb_clx]
            # Walk Clx entries; Prc (0x01) is skipped, Pcdt (0x02) holds the piece table
            i = 0
            pcdt = None
            while i < len(clx):
                if clx[i] == 0x01:
                    cb = struct.unpack_from("<H", clx, i + 1)[0]
                    i += 3 + cb
                elif clx[i] == 0x02:
                    lcb = struct.unpack_from("<I", clx, i + 1)[0]
                    pcdt = clx[i + 5: i + 5 + lcb]
                    break
                else:
                    break

            if pcdt:
                # PlcPcd: (n+1) character-position ints, then n 8-byte PCDs
                n = (len(pcdt) - 4) // 12
                cps = [struct.unpack_from("<I", pcdt, k * 4)[0] for k in range(n + 1)]
                pcd_off = (n + 1) * 4
                parts = []
                for k in range(n):
                    fc_field = struct.unpack_from("<I", pcdt, pcd_off + k * 8 + 2)[0]
                    compressed = bool(fc_field & 0x40000000)
                    fc = fc_field & 0x3FFFFFFF
                    nchars = cps[k + 1] - cps[k]
                    if nchars <= 0:
                        continue
                    if compressed:
                        off = fc // 2
                        parts.append(wd[off: off + nchars].decode("cp1255", errors="ignore"))
                    else:
                        parts.append(wd[fc: fc + nchars * 2].decode("utf-16-le", errors="ignore"))
                text = "".join(parts)

        # Fallback for documents without a usable piece table: contiguous fcMin region
        if len(text.strip()) < 5:
            fc_min = struct.unpack_from("<I", wd, 0x18)[0]
            ccp = struct.unpack_from("<I", wd, 0x4C)[0]
            if 0 < ccp and fc_min < len(wd):
                text = wd[fc_min: fc_min + ccp * 2].decode("utf-16-le", errors="ignore")

        cleaned = _clean_ole_text(text)
        return cleaned, len(cleaned) > 5
    except Exception as e:
        try:
            ole.close()
        except Exception:
            pass
        return f"[.doc extraction error: {e}]", False


def _extract_ppt_ole(path: str) -> tuple[str, bool]:
    """Pure-Python binary .ppt extraction.

    PowerPoint records are nested: container records (recVer == 0xF) hold child
    records. A flat scan misses text atoms buried inside containers, so this walks
    the record tree recursively, collecting TextCharsAtom (0x0FA0, UTF-16LE) and
    TextBytesAtom (0x0FA8, 8-bit) payloads.
    """
    import struct
    try:
        import olefile
    except ImportError:
        return "[.ppt: olefile not installed]", False
    try:
        ole = olefile.OleFileIO(path)
    except Exception as e:
        return f"[.ppt: not a valid OLE file: {e}]", False

    try:
        if not ole.exists("PowerPoint Document"):
            ole.close()
            return "[.ppt: no PowerPoint Document stream]", False
        data = ole.openstream("PowerPoint Document").read()
        ole.close()

        texts: list[str] = []

        def walk(start: int, end: int):
            i = start
            while i + 8 <= end:
                ver_inst = struct.unpack_from("<H", data, i)[0]
                rec_type = struct.unpack_from("<H", data, i + 2)[0]
                rec_len = struct.unpack_from("<I", data, i + 4)[0]
                body_start = i + 8
                body_end = body_start + rec_len
                if body_end > end:
                    break
                if (ver_inst & 0x000F) == 0x0F:      # container — descend
                    walk(body_start, body_end)
                elif rec_type == 0x0FA0:              # TextCharsAtom (UTF-16LE)
                    texts.append(data[body_start:body_end].decode("utf-16-le", errors="ignore"))
                elif rec_type == 0x0FA8:              # TextBytesAtom (8-bit)
                    texts.append(data[body_start:body_end].decode("cp1255", errors="ignore"))
                i = body_end

        walk(0, len(data))
        result = _clean_ole_text("\n".join(t for t in texts if t.strip()))
        return result, len(result) > 5
    except Exception as e:
        return f"[.ppt extraction error: {e}]", False


def extract_text(path: str) -> tuple[str, bool]:
    """Return (text, success). Imports deferred so missing libs don't crash startup."""
    ext = os.path.splitext(path)[1].lower()
    path = long_path(path)  # Windows extended-length so >260-char source paths open
    converted_tmp = None
    try:
        if ext == ".xlsx":
            import pandas as pd
            dfs = pd.read_excel(path, sheet_name=None, engine="openpyxl")
            parts = []
            for sheet, df in dfs.items():
                # Excel "used range" often extends thousands of rows/cols past the
                # actual data (stray formatting, cleared cells). Dropping all-empty
                # rows and columns prevents a tiny sheet from serializing into MBs
                # of tab separators.
                df = df.dropna(axis=0, how="all").dropna(axis=1, how="all")
                if df.empty:
                    continue
                parts.append(f"[Sheet: {sheet}]\n" + df.to_csv(sep="\t", index=False))
            return "\n".join(parts), True

        elif ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(path)
            parts = []
            for sheet in wb.sheets():
                grid = [[str(sheet.cell_value(r, c)).strip() for c in range(sheet.ncols)]
                        for r in range(sheet.nrows)]
                # keep only columns that carry data in at least one row, and skip
                # rows that are entirely empty across those columns
                keep_cols = [c for c in range(sheet.ncols) if any(row[c] for row in grid)]
                if not keep_cols:
                    continue
                rows = ["\t".join(row[c] for c in keep_cols)
                        for row in grid if any(row[c] for c in keep_cols)]
                if rows:
                    parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
            return "\n".join(parts), True

        elif ext == ".doc":
            # Prefer LibreOffice if available; fall back to pure-Python OLE reader
            converted_tmp = _libreoffice_convert(path, "docx")
            if converted_tmp:
                from docx import Document
                doc = Document(converted_tmp)
                return "\n".join(p.text for p in doc.paragraphs), True
            return _extract_doc_ole(path)

        elif ext == ".docx":
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs), True

        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(path)
            return "".join(page.extract_text() or "" for page in reader.pages), True

        elif ext == ".ppt":
            # Prefer LibreOffice if available; fall back to pure-Python OLE reader
            converted_tmp = _libreoffice_convert(path, "pptx")
            if converted_tmp:
                from pptx import Presentation
                prs = Presentation(converted_tmp)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts), True
            return _extract_ppt_ole(path)

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts), True

        elif ext in (".txt", ".md", ".py", ".java", ".html", ".csv"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(), True

        return "", False
    except Exception as e:
        return f"[extraction error: {e}]", False
    finally:
        if converted_tmp and os.path.exists(converted_tmp):
            try:
                shutil.rmtree(os.path.dirname(converted_tmp), ignore_errors=True)
            except Exception:
                pass


def clean_text(text: str) -> str:
    if not text:
        return ""
    import re
    text = str(text).replace("\x00", "")
    # rstrip each line so trailing tab/space padding (e.g. sparse spreadsheet
    # rows) can't bloat the output, then collapse runs of blank lines.
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def extract_audio_from_video(video_path: str, output_mp3_path: str, ffmpeg_path: str = "ffmpeg") -> None:
    """Extract mono 16kHz mp3 from video — no video stream, ready for Whisper.

    Raises RuntimeError carrying ffmpeg's own stderr on failure, so the real cause
    is surfaced and recorded instead of collapsed into a generic message."""
    result = subprocess.run(
        [
            ffmpeg_path, "-i", long_path(video_path),
            "-vn",
            "-ac", "1",
            "-ar", "16000",
            "-acodec", "libmp3lame", "-q:a", "4",
            "-af", "aresample=async=1",
            "-y", "-hide_banner", "-loglevel", "error",
            output_mp3_path,
        ],
        capture_output=True,
    )
    if result.returncode != 0:
        err = (result.stderr or b"").decode("utf-8", "replace").strip()
        raise RuntimeError(
            f"ffmpeg audio extraction failed (exit {result.returncode}): "
            f"{err[-600:] or '(no stderr output)'}"
        )


def get_processing_type(mime_type: str, filename: str, ignore_exts=None) -> str:
    """Determine how a file should be processed."""
    ext = os.path.splitext(filename)[1].lower()

    if ignore_exts and ext in ignore_exts:
        return "skip"

    google_native = {
        "application/vnd.google-apps.document",
        "application/vnd.google-apps.spreadsheet",
        "application/vnd.google-apps.presentation",
    }
    if mime_type in google_native:
        return "google_export"

    if ext in TRANSCRIBE_EXTS:
        return "transcribe"

    if ext in SUPPORTED_EXTRACT:
        return "extract"

    return "skip"
